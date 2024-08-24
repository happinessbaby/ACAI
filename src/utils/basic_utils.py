# import requests
import os
from pathlib import Path
# import pypandoc
import uuid
import markdown
import csv
import bs4
import urllib.request
from urllib.request import Request, urlopen
import uuid
from pptx import Presentation
from typing import Any, List, Union, Dict
from docxtpl import DocxTemplate
import asyncio
from bs4 import BeautifulSoup
from playwright.async_api import async_playwright
import collections.abc, functools
from functools import wraps
from time import time
import time
import errno
import os
import signal
import functools
# import codecs
import json
import decimal
import requests
from docx import Document
import tempfile
# from odf import text, teletype
# from odf.opendocument import load
from io import BytesIO
# import docx2pdf
from pathlib import Path
# from PyPDF2 import PdfReader  
import nltk
import PyPDF2
import subprocess
import base64
import boto3
# import aspose.words as aw
import glob
from jinja2 import Template
from botocore.exceptions import ClientError
import pypandoc
import shutil
from pdf2image import convert_from_bytes
from utils.aws_manager import get_client, get_resource
import smtplib
from email.message import EmailMessage
from email.headerregistry import Address
from email.utils import make_msgid
    
from dotenv import load_dotenv, find_dotenv
from langchain_community.document_loaders import AsyncHtmlLoader, S3FileLoader, UnstructuredURLLoader
from langchain_community.document_transformers import Html2TextTransformer

_ = load_dotenv(find_dotenv()) # read local .env 
# Path to the punkt package
punkt_path = os.path.join(nltk.data.path[0], 'tokenizers', 'punkt')
# Check if 'punkt' is already downloaded
if not os.path.exists(punkt_path):
    nltk.download('punkt')



aws_access_key_id=os.environ["AWS_SERVER_PUBLIC_KEY"]
aws_secret_access_key=os.environ["AWS_SERVER_SECRET_KEY"]
user = os.getenv('USER', 'default_user')  # Get the current user or use 'default_user' if not set

STORAGE = os.environ["STORAGE"]
if STORAGE=="CLOUD":
    bucket_name = os.environ["BUCKET_NAME"]
    s3 = get_client('s3')
    base_uri = os.environ["PRODUCTION_BASE_URI"]
    libreoffice_path = os.environ["LIBREOFFICE_PATH"]
    pdftoppm_path=os.environ["PDFTOPPM_PATH"]
else:
    bucket_name=None
    s3=None
    base_uri = os.environ["BASE_URI"]
    libreoffice_path="libreoffice"
    pdftoppm_path = "pdftoppm"

def convert_to_txt(file, output_path,) -> bool:

    """ Converts file to TXT file and move it to destination location. """
    try:
        file_ext = Path(file).suffix
        if STORAGE=="LOCAL":
            if (file_ext)=='.txt' and file!=output_path:
                os.rename(file, output_path)
            elif (file_ext=='.pdf'): 
                convert_pdf_to_txt(file, output_path)
            elif (file_ext=='.docx' or file_ext==".odt"):
                pdf_path = convert_doc_to_pdf(file)
                convert_pdf_to_txt(pdf_path, output_path)
            elif (file_ext==".log"):
                convert_log_to_txt(file, output_path)
            elif (file_ext==".pptx"):
                convert_pptx_to_txt(file, output_path)
        elif STORAGE=="CLOUD":
            loader = S3FileLoader(bucket_name, file, aws_access_key_id=aws_access_key_id, aws_secret_access_key=aws_secret_access_key)
            text = loader.load()[0].page_content
            s3.put_object(Body=text, Bucket=bucket_name, Key=output_path)
            print("Successfully converted file in S3 to TXT")
        return True 
    except Exception as e:
        print(e)
        return False

def convert_doc_to_pdf(input_path, ext=".docx", max_retries=3, delay=1):
    #retrieve docx from s3
    pdf_output_path = input_path.replace(ext, '.pdf')
    output_dir = os.path.dirname(input_path)
    for attempt in range(max_retries):
        try:
            subprocess.run([libreoffice_path, '--headless', '-env:UserInstallation=file:///tmp/LibreOffice_Conversion_${user}', '--convert-to', 'pdf:writer_pdf_Export', '--outdir', output_dir, input_path], check=True)
            print('converted docx to pdf', pdf_output_path)
            return pdf_output_path
        except subprocess.CalledProcessError as e:
            print(f"Error during conversion to pdf {attempt + 1}: {e}")
            time.sleep(delay)  # Wait before retrying
    return ""  # Indicate failure after retries

def convert_pdf_to_img(pdf_path, image_format="png", max_retries=1, delay=1):
    #
    image_output_path = pdf_path.replace('.pdf', '_images')
    os.makedirs(image_output_path, exist_ok=True)
    for attempt in range(max_retries):
        try:
            # Convert PDF to images using pdftoppm
            subprocess.run([pdftoppm_path, '-{}'.format(image_format), pdf_path, image_output_path], check=True)
            # Collect the generated image paths
            image_paths = glob.glob(f"{image_output_path}-*.{image_format}")
            print("converted pdf to image: ", image_paths)
            return image_paths
        except subprocess.CalledProcessError as e:
                print(f"Error converting {pdf_path} to image on attempt {attempt + 1}: {e}")
                time.sleep(delay)  # Wait before retrying
    return []  # Indicate failure after retries



def convert_log_to_txt(file, output_path):
    with open(file, "r") as f:
        content = f.read()
        with open(output_path, "w") as f:
            f.write(content)
            f.close()
      

def convert_pptx_to_txt(pptx_file, output_path):
    prs = Presentation(pptx_file)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text+=shape.text+'\n'
    with open(output_path, 'w') as f:
        f.write(text)
        f.close()
 
#TODO: needs to find the best pdf to txt converter that takes care of special characters best (such as the dash between dates)
def convert_pdf_to_txt(pdf_file, output_path):
    # pdf = fitz.open(pdf_file)
    # pages = count_pages(pdf_file)
    # text = f"pages: {pages}"
    # for page in pdf:
    #     text+=page.get_text() + '\n'
    # with open(output_path, 'w') as f:
    #     f.write(text)
    #     f.close()# Create a PDF reader object
    pdf_file = open(pdf_file, 'rb')
    read_pdf = PyPDF2.PdfReader(pdf_file)
    pages = len(read_pdf.pages)
    text = f"pages:{pages}"
    print(text)
    for page in read_pdf.pages:
        text+=page.extract_text()
    with open(output_path, 'w') as f:
        f.write(text)
        f.close()# Create a PDF reader object



#TODO: needs to find the best docx to txt converter that takes care of special characters best
# def convert_doc_to_txt(doc_file, output_path):
#     doc = Document(doc_file)
#     # text= pypandoc.convert_file(doc_file, to="plain", format=file_ext, outputfile=output_path)
#     # print(text)
#     pages = count_pages(doc_file)
#     with open(output_path, "w") as f:
#         f.write(f"pages:{pages}")
#         for paragraph in doc.paragraphs:
#             f.write(paragraph.text + '\n')
        
# def convert_odt_to_txt(odt_file, txt_file):
#     doc = load(odt_file)
#     pages = count_pages(odt_file)
#     text_content = [f"pages:{pages}"]
#     for paragraph in doc.getElementsByType(text.P):
#         text_content.append(teletype.extractText(paragraph))
    
#     with open(txt_file, 'w', encoding='utf-8') as f:
#         for line in text_content:
#             f.write(line + '\n')

# def convert_txt_to_doc(txt_file, output_path,):
#     doc = Document()
#     with open(txt_file, 'r', encoding='utf-8') as f:
#         for line in f:
#             doc.add_paragraph(line.strip())
#     if STORAGE=="LOCAL":
#         doc.save(output_path)
#     elif STORAGE=="CLOUD":
#          # Save DOCX to memory
#         docx_bytes = BytesIO()
#         doc.save(docx_bytes)
#         # Upload DOCX to S3
#         s3.put_object(Body=docx_bytes.getvalue(), Bucket=bucket_name, Key=output_path)

def read_txt(file: str, ) -> str:

    """ Reads TXT file into string. """

    try:
        if STORAGE=="LOCAL":
            with open(file, 'r', errors='ignore') as f:
                text = f.read()
                return text
        elif STORAGE=="CLOUD":
            data = s3.get_object(Bucket=bucket_name, Key=file)
            contents = data['Body'].read()
            text = contents.decode("utf-8")
            return text
    except Exception as e:
        raise e
    
def delete_file(file, ) -> bool:
    
    """ Deletes file. """
    
    try:
        if STORAGE=="LOCAL":
            os.remove(file)
        elif STORAGE=="CLOUD":
            s3.delete_object(Bucket=bucket_name, Key=file)
        return True
    except Exception as e:
        return False
    
def mk_dirs(paths: List[str],):

    """ Creates directories recursively given a list of paths"""

    if STORAGE=="LOCAL":
        for path in paths:
            try: 
                os.makedirs(path, exist_ok=True)
                print("Successfully made directories")
            except FileExistsError:
                pass
    elif STORAGE=="CLOUD":
        for path in paths:
            parts = path.split('/')
            current_path = ''   
            for part in parts:
                if part:  # to avoid empty strings resulting from leading/trailing slashes
                    current_path = f"{current_path}{part}/"
                    try:
                        s3.put_object(Bucket=bucket_name, Body='', Key=current_path)
                    except ClientError as e:
                        error_code = e.response['Error']['Code']
                        if error_code != 'BucketAlreadyOwnedByYou':
                            print(f"Error creating {current_path}: {e}")
               


def write_file(end_path: str, file_content="", file_path="", mode="wb",):

    """ Writes content to file. """

    if STORAGE=="LOCAL":
        try:
            with open(end_path, mode) as f:
                f.write(file_content)
                return True
        except Exception as e:
            print(e)
            return False
    elif STORAGE=="CLOUD":
        try:
            if file_path:
                s3.upload_file(file_path, bucket_name, end_path)
            else:
                s3.put_object(Body=file_content, Bucket=bucket_name, Key=end_path,)
            return True
        except Exception as e:
            print(e)
            return False

def read_file(file_path:str, mode="r", ):
    
    if STORAGE=="LOCAL":
        try:
            with open(file_path, mode) as f:
                data = f.read()
        except Exception as e:
            raise e
    elif STORAGE=="CLOUD":
        try:
            object = s3.get_object(Bucket=bucket_name, Key=file_path)
            # Read the data and decode it to a string
            data = object['Body'].read().decode('utf-8')
        except Exception as e:
            raise e
    print(data)
    return data

def move_file(source_file:str, dest_dir:str, ):

    if STORAGE=="LOCAL":
        os.rename(source_file, dest_dir)
    elif STORAGE=="CLOUD":
        s3.copy_object(
            Bucket=bucket_name,
            Key=dest_dir,
            CopySource={'Bucket': bucket_name, 'Key': source_file}
        )

def count_length(filename, ):


    # Now you can use the word tokenizer
    from nltk.tokenize import word_tokenize
    try:
        content = read_file(filename, )
        words = word_tokenize(content)
        word_count = len(words)
        return word_count
    except FileNotFoundError:
        print(f"The file {filename} does not exist.")
        return 0

    
def markdown_table_to_dict(markdown_table):
    # Convert Markdown to HTML
    html = markdown.markdown(markdown_table)
    print(html)

    # Parse HTML table using csv module
    rows = csv.reader(html.split('\n'), delimiter='|')

    # Extract header and data rows
    header = next(rows)
    data = [row for row in rows if len(row) > 1]

    # Convert rows to dictionary
    result = []
    for row in data:
        item = {}
        for i, value in enumerate(row):
            key = header[i].strip()
            item[key] = value.strip()
        result.append(item)

    return result

class AppURLopener(urllib.request.FancyURLopener):
    version = "Mozilla/5.0"

def retrieve_web_content(link, save_path="test.txt"):

    req = Request(
        url=link, 
        headers={'User-Agent': 'Mozilla/5.0'}
    )
    try: 
        webpage=str(urllib.request.urlopen(link).read())
    except Exception: 
        # webpage = urlopen(req).read()
        opener = AppURLopener()
        webpage = opener.open(link)
    soup = bs4.BeautifulSoup(webpage, features="lxml")

    content = soup.get_text()
    print(content)

    
# this one is better than the above function 
def html_to_text(urls:List[str], save_path, ):

    """Writes a list of urls' content to txt file. """
    
    try:
        loader = AsyncHtmlLoader(urls)
        docs = loader.load()
        html2text = Html2TextTransformer()
        docs_transformed = html2text.transform_documents(docs)
        content = docs_transformed[0].page_content  
        if STORAGE=="LOCAL":            
            with open(save_path, 'w') as file:
                file.write(content)
                file.close()
                print('Content retrieved and written to file.')
        elif STORAGE=="S3":
            s3.put_object(Body=content, Bucket=bucket_name, Key=save_path)
        return True
    except Exception as e:
        print(e)
        return False
    
def save_website_as_html(url, filename):
    
    try:
        response = requests.get(url)
        if response.status_code == 200:
            with open(filename, 'w', encoding='utf-8') as file:
                file.write(response.text)
            print(f"Website saved as {filename} successfully.")
        else:
            print(f"Failed to retrieve the website. Status code: {response.status_code}")
    except Exception as e:
        print(f"An error occurred: {str(e)}")


def binary_file_downloader_html(file: str, text:str="Download link") -> str:

    """ Creates the download link for AI generated file. 
    
    Args: 
    
        file: file path
        
    Returns:

        a link tag that includes the href to the file location   

    """

    data = read_file(file,mode="rb")
    bin_str = base64.b64encode(data).decode() 
    href = f'<a href="data:application/octet-stream;base64,{bin_str}" download="{os.path.basename(file)}" class="general-button">{text}</a>'
    return href

def convert_docx_to_img(docx_file_path, image_format='png'):

    def extract_idx(filename):
        basename = os.path.basename(filename)
        parts = basename.split('-')
        idx_part = parts[0].split('_')[1]
        return int(idx_part)
    
    # Convert DOCX to PDF using LibreOffice
    pdf_path = convert_doc_to_pdf(docx_file_path)

    if pdf_path:
        print("pdf_path", pdf_path)
        # Convert PDF to Image using pdftoppm
        image_paths=convert_pdf_to_img(pdf_path, image_format)
    else:
        image_paths = []

    # Return path to the image
    # pattern = os.path.join(output_dir, f"image_{idx}-*.{image_format}")
    # matching_files = glob.glob(pattern)
    # # Sort the files by extracted idx
    # image_paths = sorted(matching_files, key=extract_idx)

    return image_paths, pdf_path


def list_files(root_dir, ext=""):
    
    # Initialize a list to store the file keys
    files = []
    # Use paginator to handle large number of files
    if STORAGE=="CLOUD":
        paginator = s3.get_paginator('list_objects_v2')
        for page in paginator.paginate(Bucket=bucket_name, Prefix=root_dir):
            if 'Contents' in page:
                for obj in page['Contents']:
                    key = obj['Key']
                    print(key)
                    if ext:
                        if key.endswith(ext):
                            files.append(key)
                            print(files)
                    else:
                        files.append(key)
    elif STORAGE=="LOCAL":
        for dirpath, dirnames, filenames in os.walk(root_dir):
            for filename in filenames:
                if ext:
                    if filename.endswith(ext):
                        files.append(os.path.join(dirpath, filename))
                else:
                    files.append(os.path.join(dirpath, filename))

    return files

def get_first_file_in_each_directory(directories):
    first_files = []
    for directory in directories:
        # Use glob to find files in the directory
        if directory:
            files = glob.glob(os.path.join(directory, '*'))
            # Filter out directories, keep only files
            files = [f for f in files if os.path.isfile(f)]
            if files:
                # Sort the files to get a consistent first file (glob doesn't guarantee order)
                files.sort()
                first_files.append(files[0])
    return first_files

def write_to_docx_template(doc: Any, field_name: List[str], field_content: Dict[str, str], res_path) -> None:
    context = {key: None for key in field_name}
    for field in field_name:
        if field_content[field] != -1:
            context[field] = field_content[field]
    doc.render(context)
    doc.save(res_path)
    print(f"Succesfully written {field_name} to {res_path}.")

# source code: https://github.com/trancethehuman/entities-extraction-web-scraper/blob/main/scrape.py
async def ascrape_playwright(url, tags: list[str] = ["h1", "h2", "h3"]) -> str:
    """
    An asynchronous Python function that uses Playwright to scrape
    content from a given URL, extracting specified HTML tags and removing unwanted tags and unnecessary
    lines.
    """
    print("Started scraping...")
    results = ""
    async with async_playwright() as p:
        browser = await p.chromium.launch(headless=True)
        try:
            page = await browser.new_page()
            await page.goto(url)
            page_source = await page.content()

            # results = remove_unessesary_lines(extract_tags(remove_unwanted_tags(
            #     page_source), tags))
            text = remove_unwanted_tags(
                page_source).get_text()
            lines = (line.strip() for line in text.splitlines())
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            results = "\n".join(chunk for chunk in chunks if chunk)
            print(results)
            print("Content scraped")
        except Exception as e:
            results = f"Error: {e}"
        await browser.close()
    # return page_source
    return results 


def send_recovery_email(to_email, type, subject="Recover your password/username", password=None, username=None, ):

    # Create the base text message.
    msg = EmailMessage()
    to_email = to_email.split("@")
    # print(to_email)
    msg['Subject'] = subject
    msg['From'] = Address("aCareerAi", "yueqipeng2021", "gmail.com")
    msg['To'] = (
                Address(username, "yueqipeng2021", "gmail.com"),
                Address(username, to_email[0], to_email[1])
                 )
    if type=="username":
        msg.set_content(f"""\
            Your Username is:
                        {username}

        """)
    elif type=="password":
        token = str(uuid.uuid4())
        link = base_uri+f"""/user?token={token}&username={username}"""
        msg.set_content(f"""\
            Please follow the link to reset your passsword:
                        {link}

        """)

    msg.add_alternative(f"""\
        <!DOCTYPE html>
        <html>
            <body>
                <p>Please follow the link to reset your password:</p>
                <p><a href="{link}">Reset your password</a></p>
            </body>
        </html>
        """, subtype='html')

    smtp_server = 'smtp.gmail.com'
    smtp_port = 587
    smtp_username =  os.environ["SMTP_USERNAME"]
    smtp_password = os.environ["SMTP_PASSWORD"]
    try:
        with smtplib.SMTP(smtp_server, smtp_port) as server:
            server.starttls()  # Upgrade to secure connection
            server.login(smtp_username, smtp_password)
            server.send_message(msg)
        print("Email sent successfully!")
        return True
    except Exception as e:
        print(f"Failed to send email: {e}")
        return False

def remove_unwanted_tags(html_content, unwanted_tags=["script", "style"]) -> BeautifulSoup:
    """
    This removes unwanted HTML tags from the given HTML content.
    """
    soup = BeautifulSoup(html_content, 'html.parser')

    for tag in unwanted_tags:
        for element in soup.find_all(tag):
            element.decompose()

    return soup


def extract_tags(html_content, tags: list[str]):
    """
    This takes in HTML content and a list of tags, and returns a string
    containing the text content of all elements with those tags, along with their href attribute if the
    tag is an "a" tag.
    """
    soup = BeautifulSoup(html_content, 'html.parser')
    text_parts = []

    for tag in tags:
        elements = soup.find_all(tag)
        for element in elements:
            # If the tag is a link (a tag), append its href as well
            if tag == "a":
                href = element.get('href')
                if href:
                    text_parts.append(f"{element.get_text()} ({href})")
                else:
                    text_parts.append(element.get_text())
            else:
                text_parts.append(element.get_text())

    return ' '.join(text_parts)


def remove_unessesary_lines(content):
    # Split content into lines
    lines = content.split("\n")

    # Strip whitespace for each line
    stripped_lines = [line.strip() for line in lines]

    # Filter out empty lines
    non_empty_lines = [line for line in stripped_lines if line]

    # Remove duplicated lines (while preserving order)
    seen = set()
    deduped_lines = [line for line in non_empty_lines if not (
        line in seen or seen.add(line))]

    # Join the cleaned lines without any separators (remove newlines)
    cleaned_content = " ".join(deduped_lines)

    return cleaned_content  

def process_json(json_str: str) -> str:

    """ Processes str into valid json string """

    return json_str.strip("'<>() ").replace(" ", "").__str__().replace("'", '"')

def read_text_boxes(file_path):
    doc = Document(file_path)
    text_boxes = []
    for shape in doc.inline_shapes:
        if shape.type == 1:  # 1 is the type for text boxes
            print("text box text")
            text_box_text = []
            for paragraph in shape._inline.graphic.graphicData.textBody.p:
                text_box_text.append(paragraph.text)
            text_boxes.append('\n'.join(text_box_text))
    return text_boxes

def render_template(template_str, context):
    template = Template(template_str)
    return template.render(context)

def save_rendered_content(rendered_contents, output_file_path):
    doc = Document()
    for content in rendered_contents:
        doc.add_paragraph(content)
    doc.save(output_file_path)

class memoized(object):

    '''Decorator. Caches a function's return value each time it is called.
    If called later with the same arguments, the cached value is returned
    (not reevaluated).
    '''
    def __init__(self, func):
       self.func = func
       self.cache = {}
    def __call__(self, *args):
       if not isinstance(args, collections.abc.Hashable):
          # uncacheable. a list, for instance.
          # better to not cache than blow up.
          return self.func(*args)
       if args in self.cache:
          return self.cache[args]
       else:
          value = self.func(*args)
          self.cache[args] = value
          return value
    def __repr__(self):
       '''Return the function's docstring.'''
       return self.func.__doc__
    def __get__(self, obj, objtype):
       '''Support instance methods.'''
       return functools.partial(self.__call__, obj)
    


def timing(f):
    @wraps(f)
    def wrap(*args, **kw):
        ts = time()
        result = f(*args, **kw)
        te = time()
        print('func:%r args:[%r, %r] took: %2.4f sec' % \
          (f.__name__, args, kw, te-ts))
        return result
    return wrap

class TimeoutError(Exception):
    pass

def timeout(seconds=10, error_message=os.strerror(errno.ETIME)):
    def decorator(func):
        def _handle_timeout(signum, frame):
            raise TimeoutError(error_message)

        @functools.wraps(func)
        def wrapper(*args, **kwargs):
            signal.signal(signal.SIGALRM, _handle_timeout)
            signal.alarm(seconds)
            try:
                result = func(*args, **kwargs)
            finally:
                signal.alarm(0)
            return result

        return wrapper

    return decorator

class DecimalEncoder(json.JSONEncoder):
    def default(self, o):
        if isinstance(o, decimal.Decimal):
            if abs(o) % 1 > 0:
                return float(o)
            else:
                return int(o)
        return super(DecimalEncoder, self).default(o)

if __name__=="__main__":
    # retrieve_web_content("https://python.langchain.com/docs/use_cases/summarization/",)
    # html_to_text(
    #    "https://jobs.lever.co/missiongraduates/596a25c8-5998-469b-9cb8-53b2afd1ceab",
    #     save_path =f"./my_material/data_analyst3.txt")
        # save_path = f"./web_data/{str(uuid.uuid4())}.txt")
    convert_to_txt("/home/tebblespc/GPT-Projects/ACAI/ACAI/src/my_material/bad-resume_0.pdf","/home/tebblespc/GPT-Projects/ACAI/ACAI/src/my_material/bad-resume_0.txt")
    # convert_doc_to_txt("./test_cover_letter.docx", "docx", "./test.txt")
    # count_pages("./my_material/resume2023v2.pdf")
    # convert_docx_to_img("./backend/resume_templates/functional/functional0.docx")



    

    




